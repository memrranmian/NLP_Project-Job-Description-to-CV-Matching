"""
MatchMind - Investor Pitch Deck Generator

This script creates a professional PowerPoint presentation for VCs and angel investors
based on the MatchMind project proposal. It generates all slides with charts,
diagrams, tables, and infographics automatically.

Requirements:
    pip install python-pptx matplotlib pandas

Run this script to generate 'MatchMind_Investor_Pitch.pptx'
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import matplotlib.pyplot as plt
import pandas as pd
from io import BytesIO

def add_slide_title(pres, title_text, subtitle_text=""):
    """Add a title slide with optional subtitle"""
    slide_layout = pres.slide_layouts[0]  # Title slide layout
    slide = pres.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    if slide.placeholders[1].has_text_frame:
        slide.placeholders[1].text = subtitle_text
    return slide

def add_content_slide(pres, title_text):
    """Add a content slide with title"""
    slide_layout = pres.slide_layouts[1]  # Title and Content layout
    slide = pres.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    return slide

def add_text_to_content(slide, text, left=0.5, top=1.5, width=9, height=4, font_size=18):
    """Add text to content area with flexible positioning"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(50, 50, 50)
    return text_box

def add_bullets(slide, bullets, left=0.5, top=1.5, width=9, height=4, font_size=16):
    """Add bullet points to a slide"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0
    return text_box

def add_two_column_bullets(slide, left_bullets, right_bullets, left=0.5, top=1.5, width_left=4.2, width_right=4.2):
    """Add two columns of bullet points"""
    # Left column
    left_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width_left), Inches(4))
    left_frame = left_box.text_frame
    for i, bullet in enumerate(left_bullets):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = bullet
        p.font.size = Pt(14)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(left + width_left + 0.3), Inches(top), Inches(width_right), Inches(4))
    right_frame = right_box.text_frame
    for i, bullet in enumerate(right_bullets):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = bullet
        p.font.size = Pt(14)

def add_table(slide, data, rows, cols, left, top, width, height):
    """Add a formatted table to slide"""
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            if i < rows and j < cols:
                table.cell(i, j).text = cell_text
                # Bold header row
                if i == 0:
                    table.cell(i, j).text_frame.paragraphs[0].font.bold = True
                    table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)
                else:
                    table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)
    return table

def add_chart(slide, chart_type, data_categories, data_series, left, top, width, height, title=""):
    """Add a chart to slide"""
    chart_data = CategoryChartData()
    chart_data.categories = data_categories
    for series_name, values in data_series:
        chart_data.add_series(series_name, values)
    
    chart = slide.shapes.add_chart(
        chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    ).chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    return chart

def add_shape_with_text(slide, shape_type, left, top, width, height, text, fill_color=None):
    """Add a shape with centered text"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(12)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def create_system_flow_diagram(slide):
    """Create the 5-step system flow diagram"""
    steps = [
        ("INPUT", "Job Description + Batch Resumes"),
        ("EXTRACTION", "PDF/DOCX/TXT → Clean Text"),
        ("VECTORIZATION", "TF-IDF Vector Mapping"),
        ("SCORING", "Cosine Similarity Calculation"),
        ("RANKING", "Ranked Candidate Shortlist")
    ]
    x_start = 0.3
    y_start = 1.5
    box_width = 1.6
    box_height = 0.8
    arrow_width = 0.2
    for i, (title, subtitle) in enumerate(steps):
        x = x_start + i * (box_width + 0.15)
        # Box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y_start), Inches(box_width), Inches(box_height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(52, 152, 219)
        shape.line.color.rgb = RGBColor(41, 128, 185)
        # Text
        text_frame = shape.text_frame
        text_frame.text = f"{title}\n{subtitle}"
        text_frame.paragraphs[0].font.size = Pt(9)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Arrow between boxes
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                           Inches(x + box_width + 0.02), 
                                           Inches(y_start + box_height/2 - 0.1), 
                                           Inches(arrow_width), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)

def create_competitor_matrix(slide):
    """Create competitor comparison table"""
    data = [
        ["Competitor", "Limitation", "MatchMind Edge"],
        ["Workday / Taleo", "Expensive, keyword-heavy,\nlacks score-based ranking", "SME-friendly, deterministic\nmathematical scoring"],
        ["Generic LLMs (ChatGPT)", "Hallucination risks,\nlacks batch processing", "Zero-hallucination,\npurpose-built ranking"],
        ["Manual Review", "Slow, biased, inconsistent", "Instant, objective,\ndata-driven ranking"]
    ]
    table = slide.shapes.add_table(4, 3, Inches(0.5), Inches(1.5), Inches(9), Inches(2.5)).table
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            table.cell(i, j).text = cell_text
            if i == 0:
                table.cell(i, j).text_frame.paragraphs[0].font.bold = True
                table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)
            else:
                table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(10)

def create_revenue_chart(slide):
    """Create revenue model pie chart"""
    chart_data = CategoryChartData()
    chart_data.categories = ["SaaS Subscription", "Enterprise API", "Pay-per-Scan"]
    chart_data.add_series("Projected Revenue Mix", [60, 25, 15])
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(0.5), Inches(2.5), Inches(4), Inches(3), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.chart_title.text_frame.text = "Revenue Model (Projected Mix)"
    return chart

def create_roadmap_timeline(slide):
    """Create strategic roadmap timeline"""
    milestones = [
        ("Phase 1", "Current MVP\nFlask + TF-IDF"),
        ("Phase 2", "BERT Integration\nDeep Learning Embeddings"),
        ("Phase 3", "Edge-AI\nOn-Premise Processing"),
        ("Phase 4", "RAG Framework\nContext-Aware Retrieval"),
        ("Phase 5", "Multimodal\n360° Candidate Profiles")
    ]
    x_start = 0.3
    y_start = 2.0
    milestone_width = 1.5
    milestone_height = 1.2
    
    for i, (phase, desc) in enumerate(milestones):
        x = x_start + i * (milestone_width + 0.2)
        # Milestone box
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y_start), Inches(milestone_width), Inches(milestone_height))
        if i == 0:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(46, 204, 113)  # Green for current
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue for future
        shape.line.color.rgb = RGBColor(41, 128, 185)
        
        # Phase text
        text_frame = shape.text_frame
        text_frame.text = f"{phase}\n{desc}"
        text_frame.paragraphs[0].font.size = Pt(9)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Connector lines
        if i < len(milestones) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          Inches(x + milestone_width + 0.05),
                                          Inches(y_start + milestone_height/2 - 0.1),
                                          Inches(0.15), Inches(0.2))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(100, 100, 100)

def create_performance_chart(slide):
    """Create performance targets bar chart"""
    chart_data = CategoryChartData()
    chart_data.categories = ["Human-System\nAlignment", "Processing\nThroughput", "Extraction\nFidelity"]
    chart_data.add_series("Target", [0.92, 0.95, 0.96])
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2.2), Inches(4.5), Inches(3.5), chart_data
    ).chart
    chart.has_legend = True
    chart.chart_title.text_frame.text = "Performance Targets"
    # Format y-axis as percentage
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0.8
    value_axis.maximum_scale = 1.0
    value_axis.tick_labels.number_format = '0%'
    return chart

def add_logo_placeholder(slide, left, top, width, height):
    """Add a placeholder shape for logo"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(52, 152, 219)
    shape.line.color.rgb = RGBColor(41, 128, 185)
    text_frame = shape.text_frame
    text_frame.text = "M"
    text_frame.paragraphs[0].font.size = Pt(24)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: TITLE ==========
    slide = add_slide_title(prs, "MatchMind", "AI-Powered Resume-to-Job Matching Platform\nInvestor Pitch | 2025")
    add_logo_placeholder(slide, 8.5, 0.3, 1.0, 1.0)
    
    # ========== SLIDE 2: FOUNDERS & EXPERTISE ==========
    slide = add_content_slide(prs, "Founding Team & Domain Expertise")
    add_bullets(slide, [
        "M Imran Mian – Head of Strategic and Compliance Initiatives",
        "Afshan Ramzan – Security and Core Architect", 
        "Dania Saleem – Principal AI/ML Research Scientist",
        "Hafsah Fatima – Chief Growth Officer (CGO)",
        "",
        "Collective expertise: NLP deployment, high-concurrency HR workflows, and",
        "enterprise-grade recruitment solutions"
    ], top=1.2, font_size=14)
    
    # ========== SLIDE 3: PROBLEM STATEMENT ==========
    slide = add_content_slide(prs, "The Problem")
    problems = [
        "📊 High Screening Volume: 1000+ applications per role → manual review impossible",
        "🔍 Keyword Dependency: Legacy ATS misses qualified candidates with different terminology",
        "⏱️ Delayed Hiring Decisions: Extended cycles increase costs & reduce productivity"
    ]
    add_bullets(slide, problems, top=1.2, font_size=18)
    
    # ========== SLIDE 4: SOLUTION OVERVIEW ==========
    slide = add_content_slide(prs, "Solution: MatchMind")
    add_bullets(slide, [
        "NLP-powered recruitment screening platform",
        "Automated candidate ranking by comparing resumes against job requirements",
        "Multi-format support: PDF, DOCX, TXT",
        "TF-IDF vectorization + Cosine Similarity scoring",
        "Decision-support system – reduces manual effort, not replaces recruiters"
    ], top=1.2)
    
    # ========== SLIDE 5: SYSTEM FLOW DIAGRAM ==========
    slide = add_content_slide(prs, "Technical Architecture: System Flow")
    create_system_flow_diagram(slide)
    add_bullets(slide, [
        "Input Phase → Text Extraction → NLP Feature Extraction → Mathematical Scoring → Result Presentation"
    ], top=3.5, font_size=12)
    
    # ========== SLIDE 6: TECH STACK ==========
    slide = add_content_slide(prs, "Technology Stack")
    tech_data = [
        ["Component", "Technology", "Strategic Reason"],
        ["Backend", "Python / Flask", "API-first, high-concurrency"],
        ["Frontend", "Bootstrap / HTML/CSS", "Clean UI for HR staff"],
        ["ML Libraries", "Scikit-learn", "Industry-standard vector ops"],
        ["File Parsing", "pdfplumber / docx2txt", "High-fidelity text recovery"]
    ]
    add_table(slide, tech_data, 5, 3, 0.5, 1.2, 9, 2)
    
    # ========== SLIDE 7: HOW IT WORKS ==========
    slide = add_content_slide(prs, "Matching Engine: TF-IDF + Cosine Similarity")
    add_bullets(slide, [
        "TF-IDF Vectorization: Identifies unique skill signatures and contextual weights",
        "Cosine Similarity: Calculates precise alignment between JD and CV vectors",
        "Production-ready pipeline delivers ranked candidate lists with zero latency"
    ], top=1.2)
    
    # ========== SLIDE 8: MARKET OPPORTUNITY ==========
    slide = add_content_slide(prs, "Market Opportunity")
    add_bullets(slide, [
        "Primary Market: SMEs, Boutique Staffing Agencies, Campus Recruitment",
        "Secondary Market: Enterprise HR Teams, RPO Providers",
        "Key Drivers: Remote hiring growth, AI-assisted talent acquisition, time-to-hire pressure"
    ], top=1.2, font_size=16)
    
    # ========== SLIDE 9: VALUE PROPOSITION ==========
    slide = add_content_slide(prs, "Value Proposition")
    add_bullets(slide, [
        "✅ Faster Hiring – Reduces resume review workload",
        "✅ Consistent Evaluation – Uniform scoring across all applicants",
        "✅ Reduced Recruitment Costs – Decreases recruiter effort overhead",
        "✅ Improved Candidate Discovery – Finds talent beyond keyword matches",
        "✅ Affordable Adoption – SME-friendly alternative to enterprise ATS"
    ], top=1.2, font_size=16)
    
    # ========== SLIDE 10: COMPETITOR ANALYSIS ==========
    slide = add_content_slide(prs, "Competitor Landscape")
    create_competitor_matrix(slide)
    
    # ========== SLIDE 11: REVENUE MODEL ==========
    slide = add_content_slide(prs, "Revenue Model")
    create_revenue_chart(slide)
    add_bullets(slide, [
        "SaaS Subscription – Tiered monthly plans",
        "Enterprise API – Integration into existing ERP systems",
        "Pay-per-Scan – On-demand credits for seasonal hiring spikes"
    ], left=5.0, top=2.5, width=4.5, font_size=14)
    
    # ========== SLIDE 12: STRATEGIC ROADMAP ==========
    slide = add_content_slide(prs, "Strategic Roadmap")
    create_roadmap_timeline(slide)
    add_bullets(slide, [
        "BERT Integration → Edge-AI → RAG Framework → Multimodal Intelligence"
    ], top=3.8, font_size=12)
    
    # ========== SLIDE 13: PERFORMANCE TARGETS ==========
    slide = add_content_slide(prs, "Performance Targets")
    create_performance_chart(slide)
    add_bullets(slide, [
        "Human-System Alignment: 92%+ precision target",
        "Processing Throughput: 95%+ text extraction accuracy",
        "Extraction Fidelity: Absolute accuracy on complex document layouts"
    ], left=5.5, top=2.0, width=4.2, font_size=13)
    
    # ========== SLIDE 14: INVESTMENT OPPORTUNITY ==========
    slide = add_content_slide(prs, "Investment Opportunity")
    add_bullets(slide, [
        "Seeking strategic funding for global recruitment market expansion",
        "Current MVP deployed – Flask-based web interface with batch processing",
        "Target segments: Boutique Staffing Agencies & Campus Recruiters",
        "Differentiation: Zero-hallucination mathematical scoring",
        "Equity Structure & Fundraising Goals: TBD (available upon discussion)"
    ], top=1.2, font_size=16)
    
    # Save presentation
    output_file = "MatchMind_Investor_Pitch.pptx"
    prs.save(output_file)
    print(f"Presentation saved as: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
    print("All charts, diagrams, and tables have been generated automatically.")
    return output_file

if __name__ == "__main__":
    main()